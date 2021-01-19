from spellchecker import SpellChecker
from textblob import TextBlob,Word
import re
import docx2txt
import os
import pandas as pd
from jellyfish import soundex
from pptx import Presentation
from PyPDF2 import PdfFileReader
from pandas import ExcelWriter
from openpyxl import load_workbook

for root, dirs, files in os.walk(r'C:\Users\rajashekhar.pbr\OneDrive - Accenture\MyConcerto\Insights Engine\test_files'):
    for file in files:
        org_words = []
        dic_words = []
        spell = SpellChecker(case_sensitive=True)
        def Tokenization(test_words):
            for word in test_words:
                if len(word) > 1:
                    if spell.unknown([word]):  # check whether word is unknown or not
                        if spell.unknown([Word(word).singularize()]):
                            if spell.unknown([spell.correction(word)]):
                                org_words.append(word)  # then that word flagged as org
                            else:
                                if str(word).isupper():
                                    org_words.append(word)
                                elif any(str(chr).isalpha() for chr in word) and any(str(chr).isdigit() for chr in word):
                                    org_words.append(word)
                                else:
                                    if soundex(word) == soundex(spell.correction(word)):
                                        dic_words.append(word)
                                    else:
                                        org_words.append(word)
        if file.endswith('.docx'):
            continue
            doc_text = docx2txt.process(root+'\\'+file)
            doc_text1 = re.sub('[^A-Za-z0-9]+',' ',doc_text)
            test_words = TextBlob(doc_text1).words
            Tokenization(test_words)
        elif file.endswith('.pptx'):
            prs = Presentation(root+'\\'+file)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_runs.append(run.text)
            ppt_text = re.sub('[^A-Za-z0-9]+', ' ', ' '.join(text_runs))
            test_words = TextBlob(ppt_text).words
            Tokenization(test_words)

        df1 = pd.DataFrame(data=list(set(org_words)))
        df2 = pd.DataFrame(data=list(set(dic_words)))
        df1.to_csv(r'C:\Users\rajashekhar.pbr\Desktop\tokens5.csv', mode='a', header=False, index=False)
        df2.to_csv(r'C:\Users\rajashekhar.pbr\Desktop\tokens6.csv', mode='a', header=False, index=False)

